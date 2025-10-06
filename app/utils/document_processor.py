"""
Document Processor Module
Extracts data from Word documents, PowerPoint presentations, and images
"""

from docx import Document
from pptx import Presentation
from PIL import Image
import re
from datetime import datetime

class DocumentProcessor:
    """Process various document formats and extract relevant data"""
    
    def __init__(self):
        self.keywords = {
            'events': ['فعالية', 'حدث', 'نشاط', 'event', 'activity', 'workshop', 'ورشة', 'مؤتمر', 'conference'],
            'statistics': ['عدد', 'إحصائية', 'number', 'count', 'total', 'إجمالي'],
            'participants': ['مشارك', 'حضور', 'participant', 'attendee', 'حاضر'],
            'locations': ['موقع', 'مكان', 'location', 'venue', 'منطقة'],
        }
    
    def process_word_document(self, filepath):
        """Extract data from Word document"""
        try:
            doc = Document(filepath)
            data = {
                'events': [],
                'activities': [],
                'statistics': {},
                'text_content': []
            }
            
            current_event = None
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                data['text_content'].append(text)
                
                # Check for event indicators
                if any(keyword in text.lower() for keyword in self.keywords['events']):
                    if current_event:
                        data['events'].append(current_event)
                    
                    current_event = {
                        'title': text[:100],
                        'description': text,
                        'date': self._extract_date(text),
                        'participants': self._extract_number(text, self.keywords['participants']),
                        'location': self._extract_location(text)
                    }
                elif current_event:
                    current_event['description'] += ' ' + text
                
                # Extract statistics
                stats = self._extract_statistics(text)
                if stats:
                    data['statistics'].update(stats)
            
            if current_event:
                data['events'].append(current_event)
            
            # Process tables
            for table in doc.tables:
                table_data = self._process_table(table)
                if table_data:
                    data['statistics'].update(table_data)
            
            return data
        
        except Exception as e:
            print(f"Error processing Word document: {e}")
            return {'events': [], 'activities': [], 'statistics': {}}
    
    def process_powerpoint(self, filepath):
        """Extract data from PowerPoint presentation"""
        try:
            prs = Presentation(filepath)
            data = {
                'events': [],
                'activities': [],
                'slides_content': []
            }
            
            for slide_num, slide in enumerate(prs.slides):
                slide_data = {
                    'slide_number': slide_num + 1,
                    'title': '',
                    'content': [],
                    'images': []
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            if shape.shape_type == 1:  # Title
                                slide_data['title'] = text
                            else:
                                slide_data['content'].append(text)
                            
                            # Check for events
                            if any(keyword in text.lower() for keyword in self.keywords['events']):
                                data['events'].append({
                                    'title': text[:100],
                                    'description': text,
                                    'slide_number': slide_num + 1
                                })
                
                data['slides_content'].append(slide_data)
            
            return data
        
        except Exception as e:
            print(f"Error processing PowerPoint: {e}")
            return {'events': [], 'activities': [], 'slides_content': []}
    
    def process_image(self, filepath):
        """Process image file"""
        try:
            img = Image.open(filepath)
            return {
                'path': filepath,
                'size': img.size,
                'format': img.format,
                'mode': img.mode
            }
        except Exception as e:
            print(f"Error processing image: {e}")
            return None
    
    def _extract_date(self, text):
        """Extract date from text"""
        # Try to find date patterns
        date_patterns = [
            r'\d{1,2}[-/]\d{1,2}[-/]\d{2,4}',
            r'\d{4}[-/]\d{1,2}[-/]\d{1,2}'
        ]
        
        for pattern in date_patterns:
            match = re.search(pattern, text)
            if match:
                return match.group()
        
        return None
    
    def _extract_number(self, text, keywords):
        """Extract numbers associated with keywords"""
        for keyword in keywords:
            if keyword in text.lower():
                # Look for numbers near the keyword
                numbers = re.findall(r'\d+', text)
                if numbers:
                    return int(numbers[0])
        return None
    
    def _extract_location(self, text):
        """Extract location information"""
        for keyword in self.keywords['locations']:
            if keyword in text.lower():
                # Extract text around the location keyword
                idx = text.lower().find(keyword)
                if idx != -1:
                    return text[max(0, idx-20):min(len(text), idx+50)]
        return None
    
    def _extract_statistics(self, text):
        """Extract statistical information from text"""
        stats = {}
        
        # Look for number patterns with context
        number_patterns = re.findall(r'(\w+\s+)?(\d+)\s+(\w+)', text)
        for pattern in number_patterns:
            prefix, number, suffix = pattern
            key = f"{prefix}{suffix}".strip()
            if key:
                stats[key] = int(number)
        
        return stats
    
    def _process_table(self, table):
        """Extract data from document table"""
        data = {}
        try:
            for i, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                if len(cells) >= 2:
                    key = cells[0]
                    value = cells[1]
                    if value.isdigit():
                        data[key] = int(value)
                    else:
                        data[key] = value
        except Exception as e:
            print(f"Error processing table: {e}")
        
        return data
